# -*- coding: utf-8 -*-
"""check_values.py が本物の欠陥を拾えることを確かめる。

この検査は、誤検知を減らすために条件を絞っている。

    ・表の中の ¥0 は「確認」にとどめる（規模の小さい区分では正当に起きる）
    ・「前月から0へ」は、前月が大きかったものだけ「要対応」にする
    ・率は前月比較の対象外にする（0.3% → 0.0% は正当な動き）
    ・比率の再計算で、(8/3) のような日付は除く

絞りすぎれば、何も報告しない検査になる。
**通っても意味のない検査**にしないため、わざと壊したファイルを作り、
実際に起きた事故の型それぞれを拾えることを確かめる。

    python tests/test_check_values.py

python-pptx が無い環境では飛ばす。
"""
import os
import sys
import tempfile
import unittest

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(ROOT, "scripts"))

try:
    from pptx import Presentation
    from pptx.util import Cm, Pt
    import check_values
    READY = True
except ImportError:
    READY = False


def deck(pages):
    """[(題, [本文行...], [表の行...] or None), ...] から pptx を作る。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for title, lines, table in pages:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Cm(1), Cm(1), Cm(20), Cm(1.5))
        tb.text_frame.text = title
        y = 3.0
        for ln in lines:
            box = s.shapes.add_textbox(Cm(1), Cm(y), Cm(20), Cm(1.2))
            box.text_frame.text = ln
            box.text_frame.paragraphs[0].font.size = Pt(12)
            y += 1.4
        if table:
            rows, cols = len(table), len(table[0])
            t = s.shapes.add_table(rows, cols, Cm(1), Cm(y), Cm(20), Cm(0.8 * rows)).table
            for i, r in enumerate(table):
                for j, v in enumerate(r):
                    t.cell(i, j).text = v
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(path)
    return path


def run(path, **kw):
    opts = dict(prev=None, period=None, ratio_tol=0.15, same_min=1000,
                strict_zero=False)
    opts.update(kw)
    return check_values.check(path, **opts)


def kinds(items):
    return [k for _, k, _, _ in items]


@unittest.skipUnless(READY, "python-pptx が必要")
class TestCheckValues(unittest.TestCase):

    def test_three_brand_headline_zero_is_major(self):
        """発端の事故。3ブランドの要点ページで「収益」の見出し値が ¥0。"""
        pages = [(f"{b} 今月の要点", ["来店", "12,345", "収益", "¥0", "転換率", "1.2%"], None)
                 for b in ("ブランドA", "ブランドB", "ブランドC")]
        major, minor = run(deck(pages))
        self.assertIn("そろって0", kinds(major))
        hit = [d for _, k, _, d in major if k == "そろって0"][0]
        self.assertIn("3 ページ", hit)

    def test_table_zero_is_minor_only(self):
        """表の中の ¥0 は、規模の小さい区分で正当に起きる。要対応にしない。"""
        table = [["区分", "セッション", "収益"],
                 ["区分S", "4", "¥0"],
                 ["区分P", "12,000", "¥90,000"]]
        major, minor = run(deck([("チャネル別", ["本文"], table)] * 3))
        self.assertNotIn("そろって0", kinds(major))
        self.assertNotIn("金額の0", kinds(major))
        self.assertIn("金額の0", kinds(minor))

    def test_strict_zero_promotes(self):
        table = [["区分", "収益"], ["SNS", "¥0"]]
        major, _ = run(deck([("t", [], table)]), strict_zero=True)
        self.assertIn("金額の0", kinds(major))

    def test_reversed_period(self):
        """実際に残っていた「2026/08/01〜07/31」。終了が開始より前。"""
        major, _ = run(deck([("集計", ["集計期間: 2026/08/01〜07/31"], None)]))
        self.assertIn("期間の逆転", kinds(major))

    def test_period_month_mismatch(self):
        major, _ = run(deck([("集計", ["集計期間: 2026/07/01〜07/31"], None)]),
                       period="2026-08")
        self.assertIn("期間のずれ", kinds(major))
        major, _ = run(deck([("集計", ["集計期間: 2026/08/01〜08/31"], None)]),
                       period="2026-08")
        self.assertEqual(major, [])

    def test_ratio_recompute(self):
        _, minor = run(deck([("比率", ["PB は 25.00% (¥3,000,000 / ¥12,000,000) を占める"], None)]))
        self.assertNotIn("比率の不一致", kinds(minor))
        _, minor = run(deck([("比率", ["PB は 35.00% (¥3,000,000 / ¥12,000,000) を占める"], None)]))
        self.assertIn("比率の不一致", kinds(minor))

    def test_ratio_ignores_dates(self):
        """「0.60% (8/3)」の (8/3) は8月3日。比率ではない。"""
        _, minor = run(deck([("日次", ["CVR 最高 0.60% (8/3)"], None)]))
        self.assertNotIn("比率の不一致", kinds(minor))

    def test_duplicate_money_in_table(self):
        table = [["区分", "収益"], ["区分P", "¥90,000"], ["区分O", "¥90,000"]]
        _, minor = run(deck([("表", [], table)]))
        self.assertIn("同額の重複", kinds(minor))

    def test_prev_nonzero_to_zero_is_major_when_big(self):
        prev = deck([("ブランドA 要点", ["収益", "¥10,000,000"], None)])
        cur = deck([("ブランドA 要点", ["収益", "¥0"], None)])
        major, _ = run(cur, prev=prev)
        self.assertIn("前月から0へ", kinds(major))

    def test_prev_small_to_zero_is_minor(self):
        prev = deck([("t", ["収益", "¥5,000"], None)])
        cur = deck([("t", ["収益", "¥0"], None)])
        major, minor = run(cur, prev=prev)
        self.assertNotIn("前月から0へ", kinds(major))
        self.assertIn("前月から0へ", kinds(minor))

    def test_prev_identical_count(self):
        prev = deck([("導線", ["TOP→商品詳細 98,000"], None)])
        cur = deck([("導線", ["TOP→商品詳細 98,000"], None)])
        _, minor = run(cur, prev=prev)
        self.assertIn("前月と同一", kinds(minor))

    def test_prev_matches_nearest_position(self):
        """章ごとに同じ題のページがある。位置の近いものと比べる。"""
        prev = deck([("全体", ["来店 1,500,000"], None),
                     ("比較", ["合計 セッション 1,300,000"], None),
                     ("ブランドA", ["来店 90,000"], None),
                     ("比較", ["合計 セッション 80,000"], None)])
        cur = deck([("全体", ["来店 1,600,000"], None),
                    ("比較", ["合計 セッション 1,400,000"], None),
                    ("ブランドA", ["来店 95,000"], None),
                    ("比較", ["合計 セッション 82,000"], None)])
        _, minor = run(cur, prev=prev)
        # 全体の比較ページとブランドの比較ページを取り違えると「桁の飛び」が出る
        self.assertNotIn("桁の飛び", kinds(minor))

    def test_year_is_not_a_count(self):
        prev = deck([("t", ["2026年8月 と 2025年8月 の比較"], None)])
        cur = deck([("t", ["2026年8月 と 2025年8月 の比較"], None)])
        _, minor = run(cur, prev=prev)
        self.assertNotIn("前月と同一", kinds(minor))

    def test_clean_deck_passes(self):
        table = [["区分", "セッション", "収益"], ["区分P", "12,000", "¥90,000"],
                 ["区分N", "63,000", "¥170,000"]]
        major, minor = run(deck([("t", ["集計期間: 2026/08/01〜08/31", "収益 ¥12.3M"], table)]),
                           period="2026-08")
        self.assertEqual(major, [])
        self.assertEqual(minor, [])


if __name__ == "__main__":
    unittest.main(verbosity=2)
