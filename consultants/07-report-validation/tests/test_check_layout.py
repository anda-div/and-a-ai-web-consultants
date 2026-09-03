# -*- coding: utf-8 -*-
"""check_layout.py が本物の欠陥を拾えることを確かめる。

この検査は、誤検知を減らすために何度も条件をゆるめている。

    ・完全に内側に収まる重なりは指摘しない（見出し帯の上のラベルなど）
    ・テキストボックスは枠ではなく文字の占める範囲で見る
    ・画像の下に隠れた装飾は指摘しない（画像は不透明で上に載る）
    ・細い画像は縦横比2.0以上のときだけ指摘する
    ・画像のゆがみは切り抜きを反映して判定する

ゆるめすぎれば、何も報告しない検査になる。
**通っても意味のない検査**にしないため、わざと壊したファイルを作り、
6種類の欠陥それぞれを拾えることを確かめる。

    python tests/test_check_layout.py

python-pptx と Pillow が無い環境では飛ばす。
"""
import os
import sys
import tempfile
import unittest

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(ROOT, "scripts"))

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Cm, Pt
    import check_layout
    READY = True
except ImportError:
    READY = False

SUMMARY = "【このページの要約】"


def textbox(slide, x, y, w, h, text, size=12):
    t = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    t.text_frame.word_wrap = True
    r = t.text_frame.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(size)
    return t


def kinds(items, page=None):
    return {k for n, k, _, _ in items if page is None or n == page}


def opts_of(**kw):
    o = dict(max_gap=3.0, min_pt=6.0, min_img_w=5.0,
             summary_prefix=SUMMARY, margin=0.02)
    o.update(kw)
    return o


def run_check(path, **kw):
    """(要対応, 確認) を返す。前月比の分は run_check3 で見る。"""
    major, minor, _same, _marks = check_layout.check(path, **opts_of(**kw))
    return major, minor


def run_check3(path, baseline=None, **kw):
    """(要対応, 確認, 前月と同じ) を返す。"""
    o = opts_of(**kw)
    base = check_layout.Baseline.load(baseline, **o) if baseline else None
    major, minor, same, _marks = check_layout.check(path, baseline=base, **o)
    return major, minor, same


@unittest.skipUnless(READY, "python-pptx / Pillow が無い環境のため飛ばす")
class CheckLayoutTest(unittest.TestCase):

    @classmethod
    def setUpClass(cls):
        cls.tmp = tempfile.TemporaryDirectory()
        d = cls.tmp.name
        wide = os.path.join(d, "wide.png")
        Image.new("RGB", (1600, 900), (200, 210, 220)).save(wide)
        tall = os.path.join(d, "tall.png")
        Image.new("RGB", (400, 4000), (210, 200, 200)).save(tall)

        prs = Presentation()
        prs.slide_width = Cm(27.52)
        prs.slide_height = Cm(19.05)
        blank = prs.slide_masters[0].slide_layouts[6]

        # 1 文字どうしが本当に重なる
        s = prs.slides.add_slide(blank)
        textbox(s, 2.0, 5.0, 12.0, 1.0,
                "これは左側に置いた長めの本文で、右の文字と重なります")
        textbox(s, 8.0, 5.2, 12.0, 1.0, "こちらが右側の本文。上の行と交差しています")
        textbox(s, 1.15, 16.6, 25.2, 1.0, SUMMARY + " 重なり")

        # 2 枠外へはみ出す
        s = prs.slides.add_slide(blank)
        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(22.0), Cm(4.0),
                               Cm(9.0), Cm(3.0))
        b.text_frame.paragraphs[0].add_run().text = "右にはみ出した箱"
        textbox(s, 1.15, 16.6, 25.2, 1.0, SUMMARY + " 枠外")

        # 3 画像を縦につぶす（切り抜きなし）
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(wide, Cm(3.0), Cm(3.0), Cm(20.0), Cm(3.0))
        textbox(s, 1.15, 16.6, 25.2, 1.0, SUMMARY + " ゆがみ")

        # 4 縦長のページ全体を1枚で貼って細くなる
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(tall, Cm(12.0), Cm(2.2), Cm(1.3), Cm(13.0))
        textbox(s, 1.15, 16.6, 25.2, 1.0, SUMMARY + " 細い画像")

        # 5 パスのエスケープ漏れ（\\r と \\b が制御文字になる）
        s = prs.slides.add_slide(blank)
        textbox(s, 2.0, 5.0, 20.0, 2.0, "> python _scripts\report\build_x.py")
        textbox(s, 1.15, 16.6, 25.2, 1.0, SUMMARY + " 制御文字")

        # 6 内容が上半分で終わる
        s = prs.slides.add_slide(blank)
        textbox(s, 2.0, 2.5, 12.0, 1.0, "内容はここで終わり")
        textbox(s, 1.15, 16.6, 25.2, 1.0, SUMMARY + " 下の空き")

        cls.path = os.path.join(d, "broken.pptx")
        prs.save(cls.path)
        cls.major, cls.minor = run_check(cls.path)

    @classmethod
    def tearDownClass(cls):
        cls.tmp.cleanup()

    # ---------------------------------------------------------- 拾えること
    def test_文字どうしの重なりを拾う(self):
        self.assertIn("重なり", kinds(self.major, 1))

    def test_枠外を拾う(self):
        self.assertIn("枠外", kinds(self.major, 2))

    def test_画像のゆがみを拾う(self):
        self.assertIn("画像のゆがみ", kinds(self.major, 3))

    def test_細い画像を拾う(self):
        self.assertIn("細い画像", kinds(self.major, 4))

    def test_制御文字を拾う(self):
        self.assertIn("制御文字", kinds(self.major, 5))

    def test_下の空きを拾う(self):
        self.assertIn("下が空いている", kinds(self.minor, 6))

    def test_要対応があれば終了コードは1(self):
        self.assertEqual(check_layout.report("broken.pptx", self.major,
                                             self.minor), 1)

    # ---------------------------------------------- 指摘してはいけないこと
    def test_見出し帯の上のラベルは指摘しない(self):
        """完全に内側に収まる重なりは、意図した重ね置きとして扱う。"""
        with tempfile.TemporaryDirectory() as d:
            prs = Presentation()
            prs.slide_width = Cm(33.87)
            prs.slide_height = Cm(19.05)
            s = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[6])
            band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0),
                                      Cm(33.87), Cm(1.98))
            band.text_frame.paragraphs[0].add_run().text = "見出し"
            textbox(s, 25.7, 0.5, 7.3, 1.0, "ラベル", size=10.5)
            p = os.path.join(d, "band.pptx")
            prs.save(p)
            major, _ = run_check(p, max_gap=99.0)
            self.assertNotIn("重なり", kinds(major))

    def test_切り抜いた画像はゆがみとしない(self):
        """右端を切り落とした画像は、元の縦横比と配置が違って当然である。"""
        with tempfile.TemporaryDirectory() as d:
            src = os.path.join(d, "w.png")
            Image.new("RGB", (2000, 1000), (220, 220, 220)).save(src)
            prs = Presentation()
            prs.slide_width = Cm(27.52)
            prs.slide_height = Cm(19.05)
            s = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[6])
            # 右半分を切り落とすと、見える部分は 1000x1000（縦横比 1.0）
            pic = s.shapes.add_picture(src, Cm(2.0), Cm(2.0), Cm(10.0), Cm(10.0))
            pic.crop_right = 0.5
            p = os.path.join(d, "crop.pptx")
            prs.save(p)
            major, _ = run_check(p, max_gap=99.0)
            self.assertNotIn("画像のゆがみ", kinds(major))

    def test_横並びに分割した画像は細さを見ない(self):
        """ヒートマップは分割して並べるため、1枚ずつは必ず細くなる。"""
        with tempfile.TemporaryDirectory() as d:
            src = os.path.join(d, "t.png")
            Image.new("RGB", (400, 4000), (210, 200, 200)).save(src)
            prs = Presentation()
            prs.slide_width = Cm(27.52)
            prs.slide_height = Cm(19.05)
            s = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[6])
            for i in range(4):
                s.shapes.add_picture(src, Cm(2.0 + i * 2.5), Cm(2.2),
                                     Cm(1.3), Cm(13.0))
            p = os.path.join(d, "split.pptx")
            prs.save(p)
            major, minor = run_check(p, max_gap=99.0)
            self.assertNotIn("細い画像", kinds(major) | kinds(minor))


@unittest.skipUnless(READY, "python-pptx / Pillow が無い環境のため飛ばす")
class BaselineTest(unittest.TestCase):
    """前月の納品ファイルを基準にする（--baseline）

    毎月同じ土台から作る資料には、図の作りそのものに由来する重なりが残る。
    それを毎月直させても意味がない。**前月より乱れたかどうか**を見る。
    ただし「一度許したら以後は見ない」ではいけないので、
    大きくなったもの・増えたもの・新しく出たものは拾えることを確かめる。
    """

    @staticmethod
    def deck(path, boxes):
        prs = Presentation()
        prs.slide_width = Cm(33.87)
        prs.slide_height = Cm(19.05)
        s = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[6])
        for x, y, w, h, text in boxes:
            sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Cm(x), Cm(y), Cm(w), Cm(h))
            sh.text_frame.paragraphs[0].add_run().text = text
            sh.name = text        # 形の名前で突き合わせるため、名前を固定する
        prs.save(path)
        return path

    # 辺が交差する重なり。内包ではないので、通常は要対応になる。
    PREV = [(2.0, 2.0, 8.0, 3.0, "帯A"), (8.0, 4.0, 8.0, 3.0, "札B")]

    def test_前月と同じ重なりは要対応にしない(self):
        with tempfile.TemporaryDirectory() as d:
            base = self.deck(os.path.join(d, "prev.pptx"), self.PREV)
            cur = self.deck(os.path.join(d, "cur.pptx"), self.PREV)
            major, _minor, same = run_check3(cur, baseline=base, max_gap=99.0)
            self.assertNotIn("重なり", kinds(major))
            self.assertIn("重なり", kinds(same))

    def test_前月より大きくなった重なりは要対応にする(self):
        with tempfile.TemporaryDirectory() as d:
            base = self.deck(os.path.join(d, "prev.pptx"), self.PREV)
            worse = [(2.0, 2.0, 8.0, 3.0, "帯A"),
                     (4.0, 3.0, 8.0, 3.0, "札B")]      # 重なりを広げる
            cur = self.deck(os.path.join(d, "cur.pptx"), worse)
            major, _minor, same = run_check3(cur, baseline=base, max_gap=99.0)
            self.assertIn("重なり", kinds(major))
            self.assertNotIn("重なり", kinds(same))

    def test_前月に無かった重なりは要対応にする(self):
        with tempfile.TemporaryDirectory() as d:
            base = self.deck(os.path.join(d, "prev.pptx"),
                             [(2.0, 2.0, 8.0, 3.0, "帯A")])
            cur = self.deck(os.path.join(d, "cur.pptx"), self.PREV)
            major, _minor, same = run_check3(cur, baseline=base, max_gap=99.0)
            self.assertIn("重なり", kinds(major))
            self.assertNotIn("重なり", kinds(same))

    def test_件数が増えたぶんは要対応にする(self):
        with tempfile.TemporaryDirectory() as d:
            base = self.deck(os.path.join(d, "prev.pptx"), self.PREV)
            more = self.PREV + [(20.0, 2.0, 8.0, 3.0, "帯A"),
                                (26.0, 4.0, 8.0, 3.0, "札B")]
            cur = self.deck(os.path.join(d, "cur.pptx"), more)
            major, _minor, same = run_check3(cur, baseline=base, max_gap=99.0)
            self.assertIn("重なり", kinds(major))
            self.assertIn("重なり", kinds(same))

    def test_前月を渡さなければこれまでどおり全件を要対応にする(self):
        with tempfile.TemporaryDirectory() as d:
            cur = self.deck(os.path.join(d, "cur.pptx"), self.PREV)
            major, _minor, same = run_check3(cur, max_gap=99.0)
            self.assertIn("重なり", kinds(major))
            self.assertEqual(same, [])


if __name__ == "__main__":
    unittest.main(verbosity=2)
