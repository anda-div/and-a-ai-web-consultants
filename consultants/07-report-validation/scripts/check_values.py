# -*- coding: utf-8 -*-
"""生成したPowerPointの「値」を検査する（納品前チェック）

体裁は check_layout.py が見る。こちらは**数字と期間が、ありえない状態になっていないか**を見る。

    python check_values.py <レポート.pptx> [--prev <前月.pptx>] [--period 2026-08]

なぜ要るのか
    レポート生成の途中で1か所直し漏れると、正しく動いているように見えて、
    値だけが違うものが出来上がる。実際に起きた例：

      ・3ブランドの収益がそろって ¥0 になっていた
        （集計側は直したが、ページを描く側が古い参照先を読み続けていた）
      ・「集計期間: 2026/08/01〜07/31」── 終了日が開始日より前
        （前月の表記の月だけを書き換えた残骸。3ページに残っていた）

    どちらも、人が全ページを開いて読めば分かる。だが読み落とす。
    **ありえない値は、人より機械のほうが確実に見つける。**

検査するもの（すべて、案件に依らない）

  要対応 ── ほぼ確実に誤り
    そろって0      見出しの値（表の中でないもの）が3ページ以上で0。
                    ブランド別・チャネル別ページの直し漏れの型
    期間の逆転      「開始〜終了」の終了が開始より前
    期間のずれ      --period を渡したとき、期間表記の開始月が対象月と違う
    前月から0へ     --prev を渡したとき、前月は値があったのに今月0

  確認 ── 正当なこともあるが、目を通す価値がある
    金額の0        ¥0 / 0円 がある（規模の小さいチャネルでは正当に起きる）
    同額の重複      同じページの別の行に、1円まで同じ金額（コピーの疑い）
    比率の再計算    「28.0% (a / b)」の a/b を計算し直して合わない
    率の範囲        率・CVR・直帰率と書かれた値が 100% を超えている
    前月と同一      --prev を渡したとき、金額・件数が前月と1円・1件も違わない
                    （更新が走っていない疑い。率は対象外）
    桁の飛び        --prev を渡したとき、前月から10倍以上／10分の1以下

  --strict-zero を付けると、金額の0 も要対応になる。

見つからないもの（人の目・AIの目が必要）

    値が正しいが解釈が違う、言い過ぎ、指標の定義違い。
    それは /final-check の Phase 2 が見る。ここは**ありえない**だけを拾う。

判断の方針
    見逃すより、誤って指摘するほうを選ぶ。ただし指摘は人が数十秒で
    流し読みできる件数に収める。そのために「要対応」と「確認」を分けた。

    実際のレポートで試したところ、表の中の ¥0 は「SNS経由 4セッション」の
    ような正当なものが大半だった。一方で、3ブランドの要点ページがそろって
    ¥0 になった事故は、表の外の見出し値だった。この違いで振り分けている。

終了コードは、要対応が1件でもあれば 1、無ければ 0。
"""
from __future__ import annotations

import argparse
import difflib
import os
import re
import sys
from collections import defaultdict

from pptx import Presentation

# ---------------------------------------------------------------- 取り出し

# 金額。¥12.3M / ¥12,345,678 / 98,765円 / ¥0 を拾う。
MONEY = re.compile(
    r"(?:[¥￥]\s*(?P<a>\d[\d,]*(?:\.\d+)?)\s*(?P<unit>[MKmk万億]?)"
    r"|(?P<b>\d[\d,]*(?:\.\d+)?)\s*(?P<unit2>[万億]?)\s*円)")
# 件数など。カンマ区切りか3桁以上の整数。％や日付の一部は後で除く。
COUNT = re.compile(r"(?<![\d./:¥￥-])(\d{1,3}(?:,\d{3})+|\d{3,})(?![\d,.%/:])")
PERCENT = re.compile(r"(-?\d+(?:\.\d+)?)\s*%")
# 期間。 2026/08/01〜07/31 / 2026/08/01〜2026/08/31 / 2026-08-01 〜 2026-08-31
PERIOD = re.compile(
    r"(?P<y1>\d{4})[/-](?P<m1>\d{1,2})[/-](?P<d1>\d{1,2})\s*[〜~～\-–]\s*"
    r"(?:(?P<y2>\d{4})[/-])?(?P<m2>\d{1,2})[/-](?P<d2>\d{1,2})")
# 「25.00% (¥3,000,000 / ¥12,000,000)」の形。
RATIO = re.compile(
    r"(?P<pct>\d+(?:\.\d+)?)\s*%\s*[（(]\s*[¥￥]?\s*(?P<a>\d[\d,]*(?:\.\d+)?)\s*"
    r"[/／]\s*[¥￥]?\s*(?P<b>\d[\d,]*(?:\.\d+)?)\s*[）)]")

RATE_WORDS = ("率", "CVR", "CTR", "比率", "割合", "シェア")
UNIT = {"M": 1e6, "m": 1e6, "K": 1e3, "k": 1e3, "万": 1e4, "億": 1e8, "": 1}


def num(s: str) -> float:
    return float(s.replace(",", ""))


def label_before(line: str, start: int) -> str:
    """数値の直前にある見出し語。「収益 ¥0」なら「収益」。

    箇条書きの区切り（・ | /）で切った最後の部分から、末尾の数字と記号を落とす。
    括弧は切らない。「メルマガ (Email) 収益」のように、括弧の中がどの行かを示すため。
    """
    head = line[:start]
    head = re.split(r"[・•|｜/／、。]", head)[-1]
    head = re.sub(r"[\d,.%¥￥:：\s]+$", "", head).strip()
    return head[-24:]


def slide_texts(prs):
    """スライドごとに (題, 行のリスト) を返す。

    行は (文字列, 表の中か) の組。表は「行見出し 列見出し: 値」に展開する。
    表の中と外を分けるのは、0 の扱いが違うため（表の中の0は正当なことが多い）。
    """
    out = []
    for s in prs.slides:
        lines, title = [], ""
        for sh in s.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs).strip()
                    if t:
                        lines.append((t, False))
                        if not title and len(t) < 60:
                            title = t
            if getattr(sh, "has_table", False) and sh.has_table:
                rows = [[re.sub(r"\s+", " ", c.text.strip()) for c in r.cells]
                        for r in sh.table.rows]
                if len(rows) >= 2:
                    hdr = rows[0]
                    for r in rows[1:]:
                        for j, v in enumerate(r[1:], 1):
                            col = hdr[j] if j < len(hdr) else ""
                            lines.append((f"{r[0]} {col}: {v}", True))
        out.append((title, lines))
    return out


def labelled_values(lines):
    """行から (見出し, 種類, 値, 表の中か) を拾う。種類は money / count / pct。

    要点ページでは「収益」という見出しと「¥10.9M」という値が別の箱に入っている。
    行の中に見出しがないときは、直前の「数字を含まない短い行」を見出しとみなす。
    """
    got = []
    last_heading = ""
    for line, in_table in lines:
        if not in_table and not re.search(r"\d", line) and len(line) <= 20:
            last_heading = line
        # 日付の並びは件数と見誤るので、期間表記の部分は先に消す
        clean = PERIOD.sub(" ", line)
        clean = re.sub(r"\d{4}[/-]\d{1,2}[/-]\d{1,2}", " ", clean)
        clean = re.sub(r"\d{1,2}/\d{1,2}", " ", clean)

        def lab(pos):
            return label_before(clean, pos) or last_heading

        used = []
        for m in MONEY.finditer(clean):
            raw = m.group("a") or m.group("b")
            unit = m.group("unit") or m.group("unit2") or ""
            got.append((lab(m.start()), "money", num(raw) * UNIT.get(unit, 1), in_table))
            used.append((m.start(), m.end()))
        for m in PERCENT.finditer(clean):
            got.append((lab(m.start()), "pct", float(m.group(1)), in_table))
            used.append((m.start(), m.end()))
        for m in COUNT.finditer(clean):
            if any(a <= m.start() < b for a, b in used):
                continue
            raw = m.group(1)
            # 「2026」「2025」は年号。件数ではない
            if "," not in raw and len(raw) == 4 and 1900 <= int(raw) <= 2100:
                continue
            got.append((lab(m.start()), "count", num(raw), in_table))
    return got


def match_prev(title: str, i: int, n_cur: int, old) -> int | None:
    """今月の i 枚目に対応する前月のスライド番号を返す。

    題だけで選ぶと、章ごとに繰り返される定型ページ（「当月と前年同月の比較」など）を
    別の章と結びつけてしまう。題が似ている候補のうち、**位置が最も近い**ものを選ぶ。
    """
    if not title:
        return None
    titles = [t for t, _ in old]
    cands = [j for j, t in enumerate(titles)
             if t and difflib.SequenceMatcher(None, title, t).ratio() >= 0.6]
    if not cands:
        return None
    expect = i * len(old) / max(n_cur, 1)
    return min(cands, key=lambda j: abs(j - expect))


# ---------------------------------------------------------------- 検査

def looks_like_date(a: float, b: float, text: str) -> bool:
    """「0.60% (8/3)」の (8/3) は日付。比率の再計算から外す。"""
    return a <= 12 and b <= 31 and "," not in text and "¥" not in text and "￥" not in text


def check(path: str, *, prev: str | None, period: str | None,
          ratio_tol: float, same_min: float, strict_zero: bool):
    major, minor = [], []
    cur = slide_texts(Presentation(path))

    zero_pages = defaultdict(set)      # 見出し → 表の外で0だったページ

    for n, (title, lines) in enumerate(cur, 1):
        joined = "\n".join(t for t, _ in lines)

        # 期間の逆転・ずれ
        for m in PERIOD.finditer(joined):
            y1, m1, d1 = int(m["y1"]), int(m["m1"]), int(m["d1"])
            y2 = int(m["y2"]) if m["y2"] else y1
            m2, d2 = int(m["m2"]), int(m["d2"])
            if (y2, m2, d2) < (y1, m1, d1):
                major.append((n, "期間の逆転", m.group(0),
                              "終了日が開始日より前。前月の表記の月だけ書き換えた残骸の疑い"))
            elif period:
                py, pm = int(period[:4]), int(period[5:7])
                if (y1, m1) != (py, pm):
                    major.append((n, "期間のずれ", m.group(0),
                                  f"対象月 {period} と開始月が違う"))

        # 比率の再計算
        for m in RATIO.finditer(joined):
            a, b, pct = num(m["a"]), num(m["b"]), float(m["pct"])
            if b and not looks_like_date(a, b, m.group(0)):
                calc = a / b * 100
                if abs(calc - pct) > ratio_tol:
                    minor.append((n, "比率の不一致", m.group(0),
                                  f"a/b を計算すると {calc:.2f}%"))

        money_seen = defaultdict(set)  # 金額 → その金額を持つ見出し
        for label, kind, v, in_table in labelled_values(lines):
            if kind == "money":
                if v == 0:
                    bucket = major if strict_zero else minor
                    bucket.append((n, "金額の0", label,
                                   "¥0。規模の小さい区分なら正当。参照先の古さも疑う"))
                    if not in_table:
                        zero_pages[label].add(n)
                elif v >= 10000 and in_table:
                    money_seen[v].add(label)
            elif kind == "pct" and any(w in label for w in RATE_WORDS) and v > 100:
                minor.append((n, "率が100%超", f"{label} {v}%", "率として不自然"))

        # 同額の重複
        for v, labels in money_seen.items():
            if len(labels) >= 2:
                minor.append((n, "同額の重複", " ／ ".join(sorted(labels)),
                              f"別の行に1円まで同じ ¥{v:,.0f}。コピーの疑い"))

    # そろって0
    for label, pages in zero_pages.items():
        if len(pages) >= 3:
            ps = sorted(pages)
            major.append((ps[0], "そろって0", label,
                          f"見出し値が {len(ps)} ページで0 "
                          f"(P{', P'.join(map(str, ps))})。区分別ページの直し漏れの型"))

    # 前月との比較。率は対象外（0.3% → 0.0% のような正当な動きを誤検知するため）
    if prev:
        old = slide_texts(Presentation(prev))
        for n, (title, lines) in enumerate(cur, 1):
            j = match_prev(title, n - 1, len(cur), old)
            if j is None:
                continue
            cur_v = {(l, k): v for l, k, v, _ in labelled_values(lines)
                     if k in ("money", "count")}
            old_v = {(l, k): v for l, k, v, _ in labelled_values(old[j][1])
                     if k in ("money", "count")}
            for key, ov in old_v.items():
                if key not in cur_v:
                    continue
                cv = cur_v[key]
                label, kind = key
                # 「大きかったものが0になった」だけを要対応にする。
                # 前月 数千円の区分が今月 ¥0 なのは、規模からみて正当に起きる
                big = ov >= (10000 if kind == "money" else 100)
                if cv == 0 and ov > 0:
                    (major if big else minor).append(
                        (n, "前月から0へ", label,
                         f"前月 {ov:,.0f} → 今月 0（前月版 P{j + 1}）"))
                elif ov == cv and abs(ov) >= same_min:
                    minor.append((n, "前月と同一", label,
                                  f"{cv:,.0f} が前月と1つも違わない。更新が走っていない疑い"))
                elif ov >= 100 and cv and (cv / ov >= 10 or cv / ov <= 0.1):
                    minor.append((n, "桁の飛び", label,
                                  f"前月 {ov:,.0f} → 今月 {cv:,.0f}（前月版 P{j + 1}）"))

    # 同じ指摘の重複を落とす
    def uniq(items):
        seen, out = set(), []
        for it in items:
            k = (it[0], it[1], it[2])
            if k not in seen:
                seen.add(k)
                out.append(it)
        return out

    return uniq(major), uniq(minor)


def report(path: str, major: list, minor: list) -> int:
    name = os.path.basename(path)
    print(f"■ {name}")
    for title, items in (("要対応", major), ("確認", minor)):
        if not items:
            continue
        print(f"\n【{title}】{len(items)} 件")
        for n, kind, label, detail in sorted(items):
            head = f"  P{n:>3}  {kind}"
            print(f"{head:<22} {detail}")
            if label:
                print(f"{'':<22} {label}")
    if not major and not minor:
        print("  指摘はありません。")
    elif not major:
        print(f"\n要対応はありません（確認 {len(minor)} 件）。")
    print()
    return 1 if major else 0


def main() -> int:
    sys.stdout.reconfigure(encoding="utf-8")
    ap = argparse.ArgumentParser(description="PowerPointの値を検査する")
    ap.add_argument("pptx", nargs="+", help="検査するファイル")
    ap.add_argument("--prev", help="前月のレポート。渡すと前月との比較も行う")
    ap.add_argument("--period", help="対象月 2026-08。渡すと期間表記のずれも見る")
    ap.add_argument("--ratio-tol", type=float, default=0.15,
                    help="比率の再計算で許す差(pt)。既定 0.15")
    ap.add_argument("--same-min", type=float, default=1000,
                    help="「前月と同一」を疑う下限。既定 1000（小さい値は偶然一致しうる）")
    ap.add_argument("--strict-zero", action="store_true",
                    help="金額の0 をすべて要対応にする（既定は確認）")
    a = ap.parse_args()

    rc = 0
    for p in a.pptx:
        if not os.path.exists(p):
            print(f"ファイルがありません: {p}")
            rc = 1
            continue
        major, minor = check(p, prev=a.prev, period=a.period,
                             ratio_tol=a.ratio_tol, same_min=a.same_min,
                             strict_zero=a.strict_zero)
        rc |= report(p, major, minor)
    return rc


if __name__ == "__main__":
    raise SystemExit(main())
