# -*- coding: utf-8 -*-
"""deck_kit — PowerPointレポートの作図部品（汎用）

クライアント名・URL・数値・文章を一切持たない。
配色・レイアウト名・座標はすべて Config（_config/branding.json）から受け取る。

将来この部品は公開リポジトリ（JOB 07 レポート・効果検証）へそのまま移す。
**ここにクライアント固有の値を書いてはいけない。**

使い方:
    from pptx import Presentation
    from report_config import load
    from deck_kit import build_kit

    cfg = load()
    prs = Presentation(cfg.template())
    K   = build_kit(prs, cfg)

    s = K.content("ページ見出し", "サブタイトル")
    K.box(s, 1.15, 2.2, 12.0, None, "見出し", [("本文", K.GRAY, False)])
    K.summary(s, "このページの要約")
    K.note(s, "出所：…")
"""
from __future__ import annotations

import os
import unicodedata as _ud
from types import SimpleNamespace

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

_RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


# ============================================================ 数値の書式
def num(v):
    """数値なら整数化、'―'等はそのまま"""
    if v is None:
        return "―"
    if isinstance(v, (int, float)):
        return int(v) if float(v) == int(v) else round(float(v), 1)
    return str(v)


def cvr_dpct(sess_cur, sess_prev, cv_cur, cv_prev):
    """CVRの前月比は、丸めたCVR同士の比ではなく生値（CV/セッション）から算出する。
    丸め値から出すと桁落ちで数値が変わり、他行と基準が揃わなくなる。"""
    try:
        a = float(cv_cur) / float(sess_cur)
        b = float(cv_prev) / float(sess_prev)
        if b == 0:
            return "―"
        d = (a / b - 1) * 100
        return ("+" if d >= 0 else "−") + f"{abs(d):.1f}%"
    except Exception:
        return "―"


def dpct(v):
    """前月比の ▲/▼ を +/− に変換（▲がマイナスと誤読されるのを避ける）"""
    t = "―" if v is None else str(v)
    return t.replace("▲", "+").replace("▼", "−")


def fmt(v, unit=""):
    n = num(v)
    if isinstance(n, int):
        return f"{n:,}{unit}"
    if isinstance(n, float):
        return f"{n}{unit}"
    return n


def fmt1(v, unit=""):
    """率（%）用。必ず小数1桁で表示し、表内の桁を揃える"""
    n = num(v)
    if isinstance(n, (int, float)):
        return f"{float(n):.1f}{unit}"
    return str(n)


def txt_w(t, size):
    """全角=1・半角=0.55として文字列の描画幅(cm)を概算する"""
    cw = size / 72 * 2.54
    return sum(1.0 if _ud.east_asian_width(ch) in "WF" else 0.55 for ch in t) * cw


# ============================================================ 本体
def build_kit(prs, cfg, *, report_key=None, strip_template_slides: bool = True):
    """Presentation と Config を束ねて、作図関数一式を返す。

    report_key を渡すと、そのレポート用の座標上書き（reports.json の
    geometry_overrides）を反映する。渡さなければ branding.json の既定値を使う。

    返り値の属性はすべて関数または定数。呼び出し側は
        run, tb, content, section, box, table, chart, summary, note …
    をそのまま使える。
    """
    B = cfg.branding
    C = cfg.colors()
    NAVY, BLUE, GRAY = C["navy"], C["blue"], C["gray"]
    LGRAY, LIGHT, WHITE = C["lgray"], C["light"], C["white"]

    LNAME = B["layouts"]
    G = cfg.geometry(report_key)
    LEFT = G["content_left"]
    WIDTH = G["content_width"]
    SUB_TOP = G["subtitle_top"]
    SUB_MAXW = G["subtitle_max_width"]
    SUMM_TOP = G["summary_top"]
    SUMM_H = G["summary_height"]
    NOTE_TOP = G["note_top"]
    NOTE_LIMIT = G["note_limit"]

    BODY_MIN = float(B["fonts"]["body_min_pt"])
    BOX_PAD = {k: v for k, v in B["box_rules"]["pad"].items() if not k.startswith("_")}
    BOX_FIT_MARGIN = B["box_rules"]["fit_margin"]
    SW, SH = cfg.slide_size_cm()

    HM = cfg.heatmap(report_key)
    HM_DIR = cfg.path(HM["source_dir"])
    PIECE_DIR = cfg.path(HM["piece_dir"])
    os.makedirs(PIECE_DIR, exist_ok=True)

    LAY = {lo.name: lo for lo in prs.slide_masters[0].slide_layouts}

    # テンプレート付属の空スライドを、スライド追加前にまとめて除去。
    # ★一覧（sldIdLst）から外すだけでは部品(slideN.xml)がパッケージに残り、
    #   後から追加するスライドと名前が衝突して「修復が必要」なファイルになる。
    #   必ず drop_rel で関連付けごと削除する。
    if strip_template_slides:
        lst = prs.slides._sldIdLst
        for e in list(lst):
            rid = e.get(_RID)
            lst.remove(e)
            if rid:
                prs.part.drop_rel(rid)
        print(f"テンプレートの既存スライドを除去（残り {len(prs.slides)} 枚）")

    OVERFLOW = []   # 本文最小サイズで収まらない箱＝文章を短くする対象
    _auto = []      # 自動高さで作った箱（equalize() で高さを揃える）

    # ------------------------------------------------------ 文字
    def kinsoku(par):
        """日本語の禁則処理を段落に明示指定する。
        python-pptx が作る段落は禁則設定を持たないため、行頭に「。」「、」「）」「」」が
        落ちる現象が起きる。eaLnBrk＝日本語の行分割、hangingPunct＝句読点のぶら下げ。"""
        pPr = par._p.get_or_add_pPr()
        pPr.set("eaLnBrk", "1")
        pPr.set("hangingPunct", "1")
        pPr.set("latinLnBrk", "0")

    def ja(r):
        """ランの言語を日本語に指定する。これが無いとPowerPointが英語として扱い、
        段落側で禁則を有効にしても適用されない。"""
        rPr = r._r.get_or_add_rPr()
        rPr.set("lang", "ja-JP")
        rPr.set("altLang", "ja-JP")
        return r

    def run(par, text, color=None, bold=False, size=11, align=None):
        if align is not None:
            par.alignment = align
        kinsoku(par)
        r = par.add_run()
        r.text = text
        r.font.color.rgb = GRAY if color is None else color
        r.font.bold = bold
        r.font.size = Pt(size)
        ja(r)
        return r

    def tb(s, x, y, w, h, wrap=True):
        t = s.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
        t.text_frame.word_wrap = wrap
        return t

    # ------------------------------------------------------ ページの型
    def set_title(s, text):
        for ph in s.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = text
                for par in ph.text_frame.paragraphs:
                    kinsoku(par)
                    for r in par.runs:
                        r.font.size = Pt(16)
                        ja(r)
                return

    def content(title, sub):
        s = prs.slides.add_slide(LAY[LNAME["content"]])
        set_title(s, title)
        # サブタイトルは1行に収める（2行になると短い行が残り体裁が崩れる）
        size = 10.5
        for cand in (10.5, 10, 9.5, 9, 8.5, 8):
            if txt_w(sub, cand) <= SUB_MAXW:
                size = cand
                break
        else:
            size = 8
        run(tb(s, LEFT, SUB_TOP, WIDTH, 0.6).text_frame.paragraphs[0], sub, BLUE, True, size)
        return s

    def section(no, title, sub):
        """章扉。sub に改行を入れると明示改行になる（語中分断を防ぐ）"""
        s = prs.slides.add_slide(LAY[LNAME["section_header"]])
        set_title(s, f"{no}　{title}")
        for ph in s.placeholders:
            if ph.placeholder_format.idx == 1:
                tf = ph.text_frame
                tf.word_wrap = True
                lines = sub.split("\n")
                tf.text = lines[0]
                for extra in lines[1:]:
                    tf.add_paragraph().text = extra
                for par in tf.paragraphs:
                    kinsoku(par)
                    for r in par.runs:
                        r.font.size = Pt(13)
                        ja(r)
        return s

    def title_slide():
        return prs.slides.add_slide(LAY[LNAME["title_slide"]])

    def summary(s, text):
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Cm(LEFT), Cm(SUMM_TOP), Cm(WIDTH), Cm(SUMM_H))
        b.fill.solid()
        b.fill.fore_color.rgb = WHITE
        b.line.color.rgb = NAVY
        b.line.width = Pt(0.75)
        b.shadow.inherit = False
        tf = b.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Cm(0.3)
        tf.margin_top = Cm(0.05)
        tf.margin_bottom = Cm(0.05)
        run(tf.paragraphs[0], "【このページの要約】 ", NAVY, True, 10)
        run(tf.paragraphs[0], text, GRAY, False, 10)
        # 要約枠は10ptなら2行まで。3行になる長さは短縮対象として記録する
        if txt_w("【このページの要約】 " + text, 10) > (WIDTH - 0.6) * 2:
            OVERFLOW.append((len(prs.slides), "【要約】" + text[:24], 0.0, 0.0))

    def note(s, text, top=None):
        top = NOTE_TOP if top is None else top
        avail = NOTE_LIMIT - top - 0.13
        w = WIDTH - 0.6
        size = 7.5
        for cand in (7.5, 7.0, 6.5, 6.0, 5.5):
            cw = cand / 72 * 2.54
            width = sum(1.0 if _ud.east_asian_width(ch) in "WF" else 0.5
                        for ch in text) * cw
            lines = max(1, -(-width // w))
            if lines * (cand * 1.25 / 72 * 2.54) <= avail:
                size = cand
                break
        else:
            print(f"    ※注記が長すぎます（{len(text)}字）: {text[:40]}…")
            size = 5.5
        run(tb(s, LEFT, top, WIDTH, NOTE_LIMIT - top).text_frame.paragraphs[0],
            text, LGRAY, False, size)

    # ------------------------------------------------------ 箱・表・グラフ
    def equalize():
        """左右に並べた箱の高さを揃える（自動高さのままだと不揃いになるため）"""
        if len(_auto) > 1:
            hh = max(sh.height for sh in _auto)
            for sh in _auto:
                sh.height = hh
        _auto.clear()

    def est_box_h(head, lines, w_cm, size, hsize):
        """箱の内容から必要な高さ(cm)を概算する"""
        usable = w_cm - 0.6
        total = 0.30
        items = ([(head, hsize)] if head else []) + [(it[0], size) for it in lines]
        for t, sz in items:
            cw = sz / 72 * 2.54
            width = sum(1.0 if _ud.east_asian_width(ch) in "WF" else 0.5
                        for ch in t) * cw
            nl = max(1, int(-(-width // usable)))
            total += nl * (sz * 1.20 / 72 * 2.54) + 0.106
        return total

    def box(s, x, y, w, h, head, lines, fill=None, hc=None, hsize=11, size=9.5,
            anchor=None, pad="loose"):
        """箱を描く。
          h=None … 内容から高さを自動決定し、pad の割合ぶん高くする（下に余白を残す）
          h=数値 … その高さに固定。収まらなければ OVERFLOW に記録（文字は縮めない）
        文字サイズは常に本文最小サイズ以上に切り上げる。"""
        fill = LIGHT if fill is None else fill
        hc = NAVY if hc is None else hc
        size = max(size, BODY_MIN)
        hsize = max(hsize, BODY_MIN + 0.5)
        if not head and not lines:
            return s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Cm(x), Cm(y), Cm(w), Cm(h or 0.1))
        need = est_box_h(head, lines, w, size, hsize)
        auto = (h is None)
        if auto:
            room = SUMM_TOP - 0.25 - y
            h = min(need * BOX_PAD.get(pad, 1.20), room)
            if h < need:
                OVERFLOW.append((len(prs.slides), head or lines[0][0][:26],
                                 round(need, 2), round(h, 2)))
        else:
            if need > h * BOX_FIT_MARGIN:
                OVERFLOW.append((len(prs.slides),
                                 head or (lines[0][0][:26] if lines else ""),
                                 round(need, 2), round(h, 2)))
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Cm(x), Cm(y), Cm(w), Cm(h))
        if auto:
            _auto.append(b)
        b.fill.solid()
        b.fill.fore_color.rgb = fill
        b.line.fill.background()
        b.shadow.inherit = False
        tf = b.text_frame
        tf.word_wrap = True
        if anchor is not None:
            tf.vertical_anchor = anchor
        tf.margin_left = Cm(0.3)
        tf.margin_right = Cm(0.3)
        tf.margin_top = Cm(0.15)
        if head:
            run(tf.paragraphs[0], head, hc, True, hsize)
        first = not head
        for item in lines:
            t, c, bd = item if len(item) == 3 else (item[0], GRAY, False)
            par = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            par.space_before = Pt(3)
            run(par, t, c, bd, size)
        return b

    def table(s, x, y, w, data, colw=None, fs=8.5, hfs=8.5, rowh=0.52,
              hcolor=None, align=None, emph=None):
        hcolor = NAVY if hcolor is None else hcolor
        nr, nc = len(data), len(data[0])
        g = s.shapes.add_table(nr, nc, Cm(x), Cm(y), Cm(w), Cm(rowh * nr)).table
        if colw:
            tot = sum(colw)
            for i, cw in enumerate(colw):
                g.columns[i].width = Cm(w * cw / tot)
        for i, row in enumerate(data):
            g.rows[i].height = Cm(rowh)
            for j, val in enumerate(row):
                c = g.cell(i, j)
                c.text = ""
                c.margin_left = Cm(0.12)
                c.margin_right = Cm(0.08)
                c.margin_top = Cm(0.02)
                c.margin_bottom = Cm(0.02)
                c.vertical_anchor = MSO_ANCHOR.MIDDLE
                par = c.text_frame.paragraphs[0]
                kinsoku(par)
                if align and j < len(align):
                    par.alignment = align[j]
                elif j > 0:
                    par.alignment = PP_ALIGN.RIGHT
                col = WHITE if i == 0 else GRAY
                bold = (i == 0)
                if emph and (i, j) in emph:
                    col, bold = emph[(i, j)]
                r = par.add_run()
                r.text = str(val)
                r.font.size = Pt(hfs if i == 0 else fs)
                r.font.color.rgb = col
                r.font.bold = bold
                ja(r)
                c.fill.solid()
                c.fill.fore_color.rgb = (hcolor if i == 0
                                         else (LIGHT if i % 2 == 0 else WHITE))
        return g

    def chart(s, kind, x, y, w, h, cats, series, legend=True, fs=8, gap=60):
        cd = CategoryChartData()
        cd.categories = cats
        for nm, vals in series:
            cd.add_series(nm, vals)
        gf = s.shapes.add_chart(kind, Cm(x), Cm(y), Cm(w), Cm(h), cd)
        ch = gf.chart
        ch.font.size = Pt(fs)
        ch.font.color.rgb = GRAY
        if legend:
            ch.has_legend = True
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch.legend.include_in_layout = False
            ch.legend.font.size = Pt(fs)
        else:
            ch.has_legend = False
        try:
            ch.value_axis.tick_labels.font.size = Pt(fs)
            ch.category_axis.tick_labels.font.size = Pt(fs)
            ch.value_axis.has_major_gridlines = True
        except Exception:
            pass
        try:
            if kind == XL_CHART_TYPE.COLUMN_CLUSTERED:
                ch.plots[0].gap_width = gap
        except Exception:
            pass
        return ch

    def _est_h(lines, w_cm, size, hsize=11):
        """箱に入れた文字の必要高さ(cm)を概算"""
        u = w_cm - 0.6
        tot = 0.3
        for item in lines:
            t = item[0]
            bd = item[2] if len(item) == 3 else False
            pt = hsize if (bd and t.startswith("◆")) else size
            cw = pt / 72 * 2.54
            wd = sum(1.0 if _ud.east_asian_width(ch) in "WF" else 0.5
                     for ch in t) * cw
            tot += max(1, -(-wd // u)) * (pt * 1.25 / 72 * 2.54) + 3 / 72 * 2.54
        return tot

    def fit_size(lines, w_cm, h_cm, sizes=(12, 11.5, 11, 10.5, 10)):
        """枠に収まる最大のフォントサイズを返す。本文最小サイズは下回らない。"""
        for sz in sizes:
            if _est_h(lines, w_cm, sz, hsize=min(12.5, sz + 1)) <= h_cm:
                return sz
        return BODY_MIN

    # ------------------------------------------------------ ヒートマップ
    def hm_split(fname, n):
        """ヒートマップを縦n分割し、(パス群, 幅px, 高px) を返す"""
        from PIL import Image as _I
        _I.MAX_IMAGE_PIXELS = None
        src = os.path.join(HM_DIR, fname)
        im = _I.open(src).convert("RGB")
        w, h = im.size
        out = []
        for i in range(n):
            fp = os.path.join(PIECE_DIR, f"{fname[:-4]}_{n}_{i+1}.png")
            if not os.path.exists(fp):
                im.crop((0, int(h * i / n), w, int(h * (i + 1) / n))).save(fp)
            out.append(fp)
        return out, w, h

    def heatmap_slide(title, sub, scroll_file, click_file, n, lines, summ, note_text):
        """Scroll map と Click map を各n分割して横並びに置く（テキスト枠は右に細く）"""
        lm, tw = HM["left_margin"], HM["text_width"]
        gi, gp = HM["gap_inner"], HM["gap_group"]
        lab_top, lab_h = HM["label_top"], HM["label_height"]
        img_top, img_h = HM["image_top"], HM["image_height"]
        order = {int(k): v for k, v in HM["order_labels"].items()}
        area = SW - lm - 0.4 - tw - 0.9

        s = content(title, sub)
        _sz = fit_size(lines, tw, 14.2)
        box(s, SW - 0.9 - tw, 2.20, tw, 14.2, None, lines,
            size=_sz, hsize=min(11, _sz + 1.5))
        W = (area - 2 * (n - 1) * gi - gp) / (2 * n)
        for g, (fn, label) in enumerate([(scroll_file, "Scroll map（到達度）"),
                                         (click_file, "Click map（クリック分布）")]):
            if fn is None:
                gx = lm + g * ((n * W + (n - 1) * gi) + gp)
                b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(gx),
                                       Cm(img_top), Cm(n * W + (n - 1) * gi), Cm(6.0))
                b.fill.solid()
                b.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
                b.line.color.rgb = LGRAY
                b.line.width = Pt(0.75)
                b.shadow.inherit = False
                tf = b.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                run(tf.paragraphs[0], "（このヒートマップは取得できていません）",
                    LGRAY, True, 10, PP_ALIGN.CENTER)
                continue
            pieces, pw_px, ph_px = hm_split(fn, n)
            r = min(W / (pw_px / 37.8), img_h / ((ph_px / n) / 37.8))
            pw = (pw_px / 37.8) * r
            ph = ((ph_px / n) / 37.8) * r
            gw = n * W + (n - 1) * gi
            gx = lm + g * (gw + gp)
            tbx = s.shapes.add_textbox(Cm(gx), Cm(lab_top), Cm(gw), Cm(lab_h))
            tbx.text_frame.word_wrap = True
            par = tbx.text_frame.paragraphs[0]
            par.alignment = PP_ALIGN.CENTER
            r1 = par.add_run()
            r1.text = label
            r1.font.size = Pt(9 if n <= 2 else 8.5)
            r1.font.bold = True
            r1.font.color.rgb = NAVY
            r2 = par.add_run()
            r2.text = order.get(n, "")
            r2.font.size = Pt(7.5)
            r2.font.color.rgb = LGRAY
            for k, fp in enumerate(pieces):
                pic = s.shapes.add_picture(fp, Cm(gx + k * (W + gi) + (W - pw) / 2),
                                           Cm(img_top), Cm(pw), Cm(ph))
                pic.line.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)
                pic.line.width = Pt(0.5)
        summary(s, summ)
        note(s, note_text)
        return s

    # ------------------------------------------------------ 品質チェック
    def report_overflow():
        """本文最小サイズで収まらなかった箱を一覧表示する。
        文字を縮めるのではなく、文章を短くして対応する。"""
        if not OVERFLOW:
            print(f"\n{BODY_MIN:.0f}ptで収まらない箱はありません（文章の短縮は不要）")
            return []
        print(f"\n★ {BODY_MIN:.0f}ptで収まらない箱 {len(OVERFLOW)}件（文章を短くしてください）")
        for pageno, head, need, have in OVERFLOW:
            print(f"   P{pageno:>3}  必要{need:>5}cm / 枠{have:>5}cm  {head}")
        return OVERFLOW

    return SimpleNamespace(
        prs=prs, LAY=LAY,
        NAVY=NAVY, BLUE=BLUE, GRAY=GRAY, LGRAY=LGRAY, LIGHT=LIGHT, WHITE=WHITE,
        SW=SW, SH=SH, LEFT=LEFT, WIDTH=WIDTH,
        SUMM_TOP=SUMM_TOP, NOTE_TOP=NOTE_TOP, NOTE_LIMIT=NOTE_LIMIT,
        BODY_MIN=BODY_MIN, OVERFLOW=OVERFLOW,
        kinsoku=kinsoku, ja=ja, run=run, tb=tb, set_title=set_title,
        content=content, section=section, title_slide=title_slide,
        summary=summary, note=note, equalize=equalize,
        est_box_h=est_box_h, box=box, table=table, chart=chart,
        fit_size=fit_size, hm_split=hm_split, heatmap_slide=heatmap_slide,
        report_overflow=report_overflow,
        num=num, fmt=fmt, fmt1=fmt1, dpct=dpct, cvr_dpct=cvr_dpct, txt_w=txt_w,
    )
