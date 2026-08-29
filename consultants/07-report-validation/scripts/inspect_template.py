# -*- coding: utf-8 -*-
"""テンプレートを解析し、branding.json に書き写す値を出力する。

クライアントからPowerPointテンプレートを受け取ったら、まずこれを実行する。
出力をそのまま branding.json へ貼れる形にしてある。

    python scripts/inspect_template.py path/to/template.pptx
    python scripts/inspect_template.py path/to/template.pptx --json
"""
from __future__ import annotations

import argparse
import collections
import json
import sys

from pptx import Presentation
from pptx.util import Emu

EMU_PER_CM = 360000


def cm(v):
    return round(v / EMU_PER_CM, 2)


def collect_colors(prs):
    """既存スライドとレイアウトから、使われている色を頻度順に集める。"""
    c = collections.Counter()
    for coll in (prs.slides, prs.slide_masters[0].slide_layouts):
        for sl in coll:
            for sh in sl.shapes:
                try:
                    if sh.has_text_frame:
                        for p in sh.text_frame.paragraphs:
                            for r in p.runs:
                                if r.font.color and r.font.color.type is not None:
                                    c[str(r.font.color.rgb)] += 1
                except Exception:
                    pass
                try:
                    if sh.fill.type == 1:
                        c[str(sh.fill.fore_color.rgb)] += 1
                except Exception:
                    pass
    return c


def collect_fonts(prs):
    c = collections.Counter()
    for sl in prs.slides:
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name:
                        c[r.font.name] += 1
    return c


def guess_geometry(prs):
    """本文の左端と幅を、既存スライドの図形の並びから推測する。"""
    lefts = collections.Counter()
    widths = collections.Counter()
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.left is None or sh.width is None:
                continue
            if sh.is_placeholder and sh.placeholder_format.idx == 0:
                continue
            lefts[cm(sh.left)] += 1
            widths[cm(sh.width)] += 1
    return lefts, widths


def main() -> int:
    ap = argparse.ArgumentParser(description="PowerPointテンプレートの解析")
    ap.add_argument("pptx")
    ap.add_argument("--json", action="store_true", help="branding.json 断片だけ出力")
    a = ap.parse_args()

    prs = Presentation(a.pptx)
    w, h = cm(prs.slide_width), cm(prs.slide_height)
    ratio = prs.slide_width / prs.slide_height
    shape = "16:9" if abs(ratio - 16 / 9) < 0.05 else (
        "4:3" if abs(ratio - 4 / 3) < 0.05 else f"{ratio:.2f}:1")
    layouts = [lo.name for lo in prs.slide_masters[0].slide_layouts]
    colors = collect_colors(prs)
    fonts = collect_fonts(prs)

    fragment = {
        "template": {"path": a.pptx},
        "slide_size_cm": {"width": w, "height": h},
        "fonts": {"japanese": "（確認して記入）", "latin": "（確認して記入）",
                  "body_min_pt": 10},
        "colors": {f"color{i + 1}": v for i, (v, _) in enumerate(colors.most_common(9))},
        "layouts": {"title_slide": "（下の一覧から選ぶ）",
                    "content": "（下の一覧から選ぶ）",
                    "section_header": "（下の一覧から選ぶ）"},
    }

    if a.json:
        print(json.dumps(fragment, ensure_ascii=False, indent=2))
        return 0

    print("=" * 66)
    print(f"テンプレート: {a.pptx}")
    print("=" * 66)
    print(f"\n■ 用紙サイズ  {w} × {h} cm　（{shape}）")
    print(f"■ 既存スライド {len(prs.slides)} 枚 / レイアウト {len(layouts)} 種")

    print("\n■ レイアウト名（branding.json の layouts に書き写す）")
    for i, n in enumerate(layouts):
        print(f"    [{i}] {n}")

    print("\n■ よく使われている色（colors に書き写す）")
    for v, n in colors.most_common(12):
        print(f"    #{v}  … {n}回")

    print("\n■ フォント（fonts に書き写す）")
    if fonts:
        for v, n in fonts.most_common(6):
            print(f"    {v}  … {n}回")
    else:
        print("    明示指定なし（テーマ既定を使用）。PowerPointで確認して記入する。")

    lefts, widths = guess_geometry(prs)
    if lefts:
        print("\n■ 本文領域の推測（geometry_cm の content_left / content_width）")
        print("    左端:", "  ".join(f"{v}cm({n})" for v, n in lefts.most_common(4)))
        print("    幅  :", "  ".join(f"{v}cm({n})" for v, n in widths.most_common(4)))
        print("    ※ 要約・注記の位置は、既存スライドを開いて実測してください。")
    else:
        print("\n■ 本文領域：既存スライドが無いため推測できません。"
              "1枚作って実測するか、既定値から始めてください。")

    print("\n■ branding.json への断片")
    print(json.dumps(fragment, ensure_ascii=False, indent=2))
    print("\n※ 「（確認して記入）」の箇所は、テンプレートをPowerPointで開いて確認してください。")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
