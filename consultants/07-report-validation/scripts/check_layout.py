# -*- coding: utf-8 -*-
"""生成したPowerPointの体裁を検査する（納品前チェック）

スライドの座標はファイルの中に正確に入っている。目で見るより速く、正確に
測れることは測ってしまい、人の目は「文章として正しいか」に使う。

    python check_layout.py <レポート.pptx> [オプション]

検査するもの

    重なり      内容どうしの辺が交差していないか
                （完全に内側に収まっているものは意図した重ね置きとみなす）
    枠外        図形が用紙からはみ出していないか
    下の空き    内容の下端と要約枠の間が空きすぎていないか
    画像のゆがみ 画像が元の縦横比のまま置かれているか
    細い画像    縦長すぎて何が写っているか読めなくなっていないか
    文字サイズ  本文が読める大きさを下回っていないか
    制御文字    パスのエスケープ漏れで _x000D_ などが混ざっていないか

見つからないもの（人の目が必要）

    解説と画面が食い違っていないか、数値の意味が正しいか、
    言い過ぎていないか、宛先・日付が正しいか。

前月の納品ファイルを基準にする（--baseline）

    毎月同じ土台から作る資料には、**図の作りそのものに由来する重なり**が
    残る。たとえばフロー図の矢印ラベルは、上下いっぱいに伸びた列の上に
    置くのが図の作りで、どこへ置いても列と辺が交差する。位置を直すと図が壊れる。

    そういうものを毎月直させても意味がない。そこで

        python check_layout.py 当月.pptx --baseline "納品)前月.pptx"

    とすると、**前月の納品ファイルに同じ形で在ったものは要対応にしない。**
    「前月と同じ」として理由付きで一覧に出し、記録には残す。

    見るのは**前月より乱れたかどうか**である。同じ図形の組み合わせでも
    重なりが前月より大きくなっていれば、要対応として拾う。
    件数が増えていても拾う。前月に無かったものは、当然拾う。

    前月のファイルを渡さなければ、これまでどおり全件を要対応にする。

終了コードは、重大な指摘が1件でもあれば 1、無ければ 0。
"""
from __future__ import annotations

import argparse
import os
import re
import sys
import unicodedata as _ud
import zipfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

EMU_CM = 360000.0


def text_width_cm(text: str, size_pt: float) -> float:
    """文字列の見た目の幅（cm）。全角は1文字、半角は0.5文字ぶんとして数える。"""
    cw = size_pt / 72 * 2.54
    return sum(1.0 if _ud.east_asian_width(c) in "WF" else 0.5
               for c in text) * cw

# 内容として扱う図形（装飾の線や矢印は対象外）
CONTENT_TYPES = {"autoShape", "pic", "graphicFrame", "tbl", "sp"}


def cm(v) -> float:
    return (v or 0) / EMU_CM


class Shape:
    """検査に必要な情報だけを取り出した図形"""

    def __init__(self, sh, slide_no: int, z: int = 0):
        self.raw = sh
        self.slide = slide_no
        self.z = z          # 描画順。大きいほど上に載る
        self.x = cm(sh.left)
        self.y = cm(sh.top)
        self.w = cm(sh.width)
        self.h = cm(sh.height)
        self.tag = sh._element.tag.rsplit("}", 1)[-1]
        self.text = ""
        if getattr(sh, "has_text_frame", False):
            self.text = sh.text_frame.text.strip()
        self.is_picture = self.tag == "pic"
        self.is_connector = self.tag == "cxnSp"
        self.is_textbox = False
        try:
            self.is_textbox = sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        except Exception:
            pass
        # テキストボックスは透明で、文字のあるところしか占めない。
        # 塗りや枠線を持つ図形は枠のぶんだけ場所を占める。
        self.tx, self.ty, self.tw, self.th = self.x, self.y, self.w, self.h
        if self.text:
            self._shrink_to_text(sh)
            # 文字の占める範囲は控えておく（画像に隠れた判定で使う）
            self.ex, self.ey, self.ew, self.eh = self.tx, self.ty, self.tw, self.th
            if not self.is_textbox:
                # 塗りや枠線を持つ図形は、枠のぶんだけ場所を占める
                self.tx, self.ty, self.tw, self.th = self.x, self.y, self.w, self.h
        else:
            self.ex, self.ey, self.ew, self.eh = self.x, self.y, self.w, self.h

    def _shrink_to_text(self, sh) -> None:
        """枠を、実際に文字が占める範囲まで縮める。"""
        usable = max(self.w - 0.2, 0.1)
        lines = 0
        widest = 0.0
        height = 0.0
        for par in sh.text_frame.paragraphs:
            t = "".join(r.text for r in par.runs)
            if not t.strip():
                lines += 1
                height += 0.35
                continue
            size = next((r.font.size.pt for r in par.runs if r.font.size), 11.0)
            w = text_width_cm(t, size)
            n = max(1, int(-(-w // usable)))
            widest = max(widest, min(w, usable))
            height += n * (size * 1.35 / 72 * 2.54)
            lines += n
        if not lines:
            return
        align = None
        try:
            align = sh.text_frame.paragraphs[0].alignment
        except Exception:
            pass
        self.tw = min(self.w, widest + 0.25)
        self.th = min(self.h, height + 0.15)
        if align == PP_ALIGN.CENTER:
            self.tx = self.x + (self.w - self.tw) / 2
        elif align == PP_ALIGN.RIGHT:
            self.tx = self.x + (self.w - self.tw)

    @property
    def right(self) -> float:
        return self.x + self.w

    @property
    def bottom(self) -> float:
        return self.y + self.h

    @property
    def label(self) -> str:
        if self.text:
            return self.text.split("\n")[0][:30]
        if self.is_picture:
            return f"画像（左 {self.x:.1f}cm・上 {self.y:.1f}cm）"
        return {"tbl": "表", "graphicFrame": "表またはグラフ"}.get(self.tag, "図形")

    @property
    def has_content(self) -> bool:
        """中身のある図形か。空の装飾用の箱は重なり判定から外す。"""
        if self.is_connector:
            return False
        if self.tag in ("pic", "graphicFrame"):
            return True
        return bool(self.text)


def collect(slide, slide_no: int) -> list[Shape]:
    """グループは中身に展開して集める（枠だけの重なりを数えないため）"""
    out = []

    def walk(shapes):
        for sh in shapes:
            if sh.shape_type is not None and sh._element.tag.endswith("}grpSp"):
                walk(sh.shapes)
                continue
            try:
                if sh.left is None or sh.top is None:
                    continue
            except (AttributeError, TypeError):
                continue
            out.append(Shape(sh, slide_no, len(out)))

    walk(slide.shapes)
    return out


def overlap_area(a: Shape, b: Shape) -> tuple[float, float, float]:
    """重なりの面積・幅・高さ。テキストボックスは文字の占める範囲で見る。"""
    ar, ab = a.tx + a.tw, a.ty + a.th
    br, bb = b.tx + b.tw, b.ty + b.th
    ow = min(ar, br) - max(a.tx, b.tx)
    oh = min(ab, bb) - max(a.ty, b.ty)
    if ow <= 0 or oh <= 0:
        return 0.0, 0.0, 0.0
    return ow * oh, ow, oh


def inside(inner: Shape, outer: Shape, tol: float = 0.05) -> bool:
    """inner が outer の内側に完全に収まっているか。

    完全な内包は、ほぼ意図した重ね置きである。
      ・見出し帯の上にラベルを置く
      ・キャプチャの上に注記を置く（改善案のBefore/After）
    事故として起きる重なりは、辺が交差する形になる。
    """
    return (inner.tx >= outer.tx - tol and inner.ty >= outer.ty - tol
            and inner.tx + inner.tw <= outer.tx + outer.tw + tol
            and inner.ty + inner.th <= outer.ty + outer.th + tol)


def control_chars(path: str) -> list[tuple[int, str, str, str]]:
    """XMLに退避された制御文字を探す。

    Windowsのパスを生の文字列にせずに書くと、\\r や \\b が
    エスケープとして解釈される。python-pptx では普通の文字列に
    見えるが、PowerPointでは _x000D_ のように表示される。
    """
    out = []
    with zipfile.ZipFile(path) as z:
        names = sorted((n for n in z.namelist()
                        if n.startswith("ppt/slides/slide")),
                       key=lambda n: int(re.findall(r"(\d+)", n)[0]))
        for n in names:
            xml = z.read(n).decode("utf-8", "ignore")
            no = int(re.findall(r"(\d+)", n)[0])
            for m in re.finditer(r"_x00[0-9A-Fa-f]{2}_", xml):
                ctx = re.sub(r"<[^>]*>", "", xml[max(0, m.start() - 60):m.end() + 30])
                out.append((no, "制御文字", ctx.strip()[:40],
                            f"{m.group(0)} が文字として入っています。"
                            "パスを書いた文字列のエスケープ漏れです"))
    return out


class Baseline:
    """前月の納品ファイルに在った指摘。ここに在るものは要対応にしない。

    突き合わせるのは**ページ番号ではなく図形の名前**である。差し込みページで
    番号がずれても、同じ形の組み合わせなら同じものとして扱える。

    同じ組み合わせでも、前月より重なりが大きくなっていれば拾う。
    件数が前月より増えていれば、増えたぶんを拾う。
    「前月と同じなら通す」であって「一度許したら以後は見ない」ではない。
    """

    # 前月より悪くなったと見なす余裕。測り方の丸めぶんは許す。
    SLACK_RATIO = 1.10
    SLACK_ABS = 0.05

    def __init__(self, marks=None, path=None):
        self.marks = marks or {}          # 印 -> [ひどさ, ...]（大きい順）
        self.path = path
        self.used = {}                    # 印 -> 使った数

    @classmethod
    def load(cls, path, **opts):
        """前月のファイルに同じ検査を当てて、指摘の印を集める。"""
        if not path:
            return cls()
        _major, _minor, _same, marks = check(path, baseline=None, **opts)
        got = {}
        for mark, harm in marks:
            got.setdefault(mark, []).append(harm)
        for v in got.values():
            v.sort(reverse=True)
        return cls(got, path)

    def allows(self, mark, harm) -> tuple[bool, str]:
        """前月に在ったものか。在れば (True, 一言) を返す。

        「ひどさ」の単位は指摘の種類によって違う（重なりは c㎡、枠外は cm、
        画像のゆがみは縦横比の差）。数を出すのは**前月より悪くなったとき**に
        限る。同じときに数を並べると、本文の寸法と混ざって読みにくい。
        """
        had = self.marks.get(mark)
        if not had:
            return False, ""
        i = self.used.get(mark, 0)
        if i >= len(had):
            return False, f"前月は同じものが {len(had)} 件、当月はそれより多い"
        was = had[i]
        if harm > was * self.SLACK_RATIO + self.SLACK_ABS:
            return False, (f"前月より大きくなった: ひどさ {was:.2f} → {harm:.2f}")
        self.used[mark] = i + 1
        return True, "前月にも同じ形の指摘がある"


def check(path: str, *, max_gap: float, min_pt: float, min_img_w: float,
          summary_prefix: str, margin: float,
          strict: bool = False, baseline: "Baseline | None" = None
          ) -> tuple[list, list, list, list]:
    """指摘を返す。

    戻り値は (要対応, 確認, 前月と同じ, 印の一覧)。
    「印の一覧」は前月のファイルを基準にするときに使う。
    """
    prs = Presentation(path)
    SW, SH = cm(prs.slide_width), cm(prs.slide_height)
    major, minor, same, marks = [], [], [], []

    def note(n, kind, label, detail, ident, harm):
        """要対応の候補を1件記録する。前月に在ったものは「前月と同じ」へ回す。"""
        mark = (kind, ident)
        marks.append((mark, harm))
        if baseline is not None:
            ok, why = baseline.allows(mark, harm)
            if ok:
                same.append((n, kind, label, f"{detail}／{why}"))
                return
            if why:
                detail = f"{detail}（{why}）"
        major.append((n, kind, label, detail))

    # 描画するまで見えないため、最初に拾っておく
    for n, kind, label, detail in control_chars(path):
        note(n, kind, label, detail, label, 1.0)

    for n, slide in enumerate(prs.slides, 1):
        shapes = collect(slide, n)
        content = [s for s in shapes if s.has_content]

        # ---------------------------------------------------- 重なり
        for i, a in enumerate(content):
            for b in content[i + 1:]:
                area, ow, oh = overlap_area(a, b)
                if area <= 0:
                    continue
                # かすっているだけのものは数えない
                if ow < 0.15 or oh < 0.15:
                    continue
                # 完全に内側に収まっているものは、意図した重ね置きとみなす
                # （見出し帯の上のラベル、キャプチャの上の注記など）。
                # 事故として起きる重なりは、辺が交差する形になる。
                if not strict and (inside(a, b) or inside(b, a)):
                    continue
                # 画像は不透明で、上に載っている。下の装飾が隠れるのは
                # 設計どおりで見た目の問題にならない。文字が隠れるときだけ見る。
                if not strict and (a.is_picture != b.is_picture):
                    pic, other = (a, b) if a.is_picture else (b, a)
                    if other.z < pic.z:
                        if not other.text:
                            continue
                        ow2 = min(other.ex + other.ew, pic.x + pic.w) - \
                            max(other.ex, pic.x)
                        oh2 = min(other.ey + other.eh, pic.y + pic.h) - \
                            max(other.ey, pic.y)
                        if ow2 <= 0.15 or oh2 <= 0.15:
                            continue
                        ow, oh = ow2, oh2
                small = min(a.tw * a.th, b.tw * b.th)
                if small <= 0 or area / small < 0.08:
                    continue
                note(n, "重なり",
                     f"{a.label} × {b.label}",
                     f"{ow:.2f} × {oh:.2f} cm 重複",
                     frozenset((a.raw.name or "", b.raw.name or "")),
                     ow * oh)

        # ---------------------------------------------------- 枠外
        for s in shapes:
            if s.is_connector:
                continue
            if s.x < -margin or s.y < -margin or s.right > SW + margin or s.bottom > SH + margin:
                over = max(-s.x, -s.y, s.right - SW, s.bottom - SH)
                note(n, "枠外", s.label, f"{over:.2f} cm はみ出し",
                     s.raw.name or "", over)

        # ---------------------------------------------------- 画像
        pics = [s for s in shapes if s.is_picture]
        split_layout = len(pics) > 1   # 分割して横並びにしているページ
        for s in shapes:
            if not s.is_picture or s.w <= 0 or s.h <= 0:
                continue
            try:
                px, py = s.raw.image.size
            except Exception:
                px = py = 0
            if px and py:
                # 切り抜き（トリミング）を反映した縦横比と比べる。
                # 右端を切り抜いた画像は、元の縦横比と配置が違って当然である。
                cl = getattr(s.raw, "crop_left", 0) or 0
                cr = getattr(s.raw, "crop_right", 0) or 0
                ct = getattr(s.raw, "crop_top", 0) or 0
                cb = getattr(s.raw, "crop_bottom", 0) or 0
                ew = px * max(1e-6, 1 - cl - cr)
                eh = py * max(1e-6, 1 - ct - cb)
                want, got = eh / ew, s.h / s.w
                if want and abs(want - got) / want > 0.02:
                    crop = "（切り抜きあり）" if (cl or cr or ct or cb) else ""
                    note(n, "画像のゆがみ", s.label,
                         f"本来 {want:.2f} → 配置 {got:.2f}{crop}",
                         s.raw.name or "", abs(want - got) / want)
            # 幅が狭いだけでは足りない。小さなアイコンは狭くて当然である。
            # 「縦に長いページ全体を1枚で貼ったため細くなった」ものだけを拾う。
            tall = (s.h / s.w) >= 2.0
            if s.w < min_img_w and tall and not split_layout:
                detail = (f"幅 {s.w:.1f} cm（縦横比 {s.h / s.w:.1f}:1）。"
                          "解説で触れた箇所を切り出してから貼る")
                if s.w < min_img_w / 2:
                    note(n, "細い画像", s.label, detail,
                         s.raw.name or "", s.h / s.w)
                else:
                    minor.append((n, "細い画像", s.label, detail))

        # ---------------------------------------------------- 文字サイズ
        # 注記（要約枠より下の帯）は小さくてよい。本文だけを見る。
        summ0 = next((s for s in shapes if s.text.startswith(summary_prefix)), None)
        body_limit = summ0.y if summ0 else SH
        for s in shapes:
            if not getattr(s.raw, "has_text_frame", False):
                continue
            if s.y >= body_limit - 0.3:
                continue
            # 箱の中の本文だけを見る。キャプションやグラフのラベルは
            # 小さくてよい（意図した使い分けを指摘しても意味がない）。
            if s.tag != "sp" or s.raw.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                continue
            for par in s.raw.text_frame.paragraphs:
                for r in par.runs:
                    if r.font.size and r.font.size.pt < min_pt and r.text.strip():
                        minor.append((n, "文字が小さい", r.text.strip()[:26],
                                      f"{r.font.size.pt:.1f} pt"))
                        break
                else:
                    continue
                break

        # ---------------------------------------------------- 下の空き
        summ = summ0
        if summ:
            limit = summ.y
            body = [s for s in content if s is not summ and s.y < limit - 0.3]
            if body:
                gap = limit - max(s.bottom for s in body)
                if gap >= max_gap:
                    minor.append((n, "下が空いている", "",
                                  f"{gap:.1f} cm。行を増やす／要素を大きくする／解釈を足す"))

    return major, minor, same, marks


def report(path: str, major: list, minor: list, same: list | None = None) -> int:
    name = os.path.basename(path)
    print(f"■ {name}")
    for title, items in (("要対応", major), ("確認", minor),
                         ("前月と同じ", same or [])):
        if not items:
            continue
        print(f"\n【{title}】{len(items)} 件")
        for n, kind, label, detail in items:
            head = f"  P{n:>3}  {kind}"
            print(f"{head:<22} {detail}")
            if label:
                print(f"{'':<22} {label}")
    if not major and not minor and not same:
        print("  指摘はありません。")
    elif not major:
        tail = f"確認 {len(minor)} 件"
        if same:
            tail += f" / 前月と同じ {len(same)} 件"
        print(f"\n要対応はありません（{tail}）。")
    print()
    return 1 if major else 0


def main() -> int:
    sys.stdout.reconfigure(encoding="utf-8")
    ap = argparse.ArgumentParser(description="PowerPointの体裁を検査する")
    ap.add_argument("pptx", nargs="+", help="検査するファイル")
    ap.add_argument("--max-gap", type=float, default=3.0,
                    help="内容の下端と要約枠の間の許容量(cm)。既定 3.0")
    ap.add_argument("--min-pt", type=float, default=6.0,
                    help="箱の中の本文の最小サイズ(pt)。既定 6.0（これを下回ると"
                         "どんな意図でも読めない）。注記帯とキャプションは対象外")
    ap.add_argument("--min-img-width", type=float, default=5.0,
                    help="画像の最小幅(cm)。既定 5.0")
    ap.add_argument("--strict-overlap", action="store_true",
                    help="完全に内側に収まっている重なりも指摘する"
                         "（見出し帯の上のラベルなども拾うため誤検知が増える）")
    ap.add_argument("--margin", type=float, default=0.02,
                    help="枠外判定の許容量(cm)。既定 0.02")
    ap.add_argument("--summary-prefix", default="【このページの要約】",
                    help="要約枠の書き出し。下の空きの判定に使う")
    ap.add_argument("--baseline",
                    help="前月の納品ファイル。ここに同じ形で在った指摘は"
                         "要対応にせず「前月と同じ」として出す。前月より"
                         "重なりが大きくなったもの・件数が増えたもの・"
                         "前月に無かったものは、これまでどおり要対応にする")
    a = ap.parse_args()

    opts = dict(max_gap=a.max_gap, min_pt=a.min_pt,
                min_img_w=a.min_img_width,
                summary_prefix=a.summary_prefix, margin=a.margin,
                strict=a.strict_overlap)

    base = None
    if a.baseline:
        if not os.path.exists(a.baseline):
            print(f"前月のファイルがありません: {a.baseline}")
            return 1
        base = Baseline.load(a.baseline, **opts)
        n_base = sum(len(v) for v in base.marks.values())
        print(f"基準: {os.path.basename(a.baseline)}（指摘 {n_base} 件）")
        print()

    rc = 0
    for p in a.pptx:
        if not os.path.exists(p):
            print(f"ファイルがありません: {p}")
            rc = 1
            continue
        if base is not None:
            base.used = {}          # ファイルごとに数え直す
        major, minor, same, _marks = check(p, baseline=base, **opts)
        rc |= report(p, major, minor, same)
    return rc


if __name__ == "__main__":
    raise SystemExit(main())
