# -*- coding: utf-8 -*-
"""site_report_kit — サイト全体レポートのページ型（汎用）

deck_kit（作図の最小部品）の上に、実績レポートで繰り返し使われている
「ページの型」を載せたもの。クライアント名・URL・数値・文章は持たない。
座標はすべて Config（reports.json の page_geometry）から受け取る。

ページ型は5種類。

  A two_col_tables … 表を左右に並べ、下に解釈の箱
  B kpi_and_chart  … 上にKPIカード4枚、下にグラフ＋右に解釈の箱
    two_charts     … グラフを左右に並べ、下に解釈の箱
  C boxes          … 箱だけで構成（サマリー・図解・まとめ）
  D capture_right  … 左に箱、右にキャプチャ（課題詳細・改善提案）
    capture_left   … 左にキャプチャ、右に箱（競合比較）
  E heatmap        … deck_kit.heatmap_slide をそのまま使う

使い方:
    from report_config import load
    from deck_kit import build_kit
    from site_report_kit import build_site_kit

    cfg = load()
    prs = Presentation(cfg.template())
    K   = build_kit(prs, cfg, report_key="site")
    S   = build_site_kit(K, cfg, report_key="site")

    S.two_col_tables(
        "01 GA4解析　流入分析", "どこから来た人が成果に近いか",
        left=[["チャネル", "セッション"], ["自然検索", "1,234"]],
        right=[["参照元", "セッション"], ["google", "1,000"]],
        insight=("注目：", [("自然検索が最大の入口です。", None, False)]),
        summ="自然検索が最大の入口。",
        note="出所：GA4（対象期間）")
"""
from __future__ import annotations

import os
from types import SimpleNamespace

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Pt


def build_site_kit(K, cfg, *, report_key: str = "site"):
    """deck_kit の名前空間 K を受け取り、ページ型の関数一式を返す。"""
    PG = cfg.page_geometry(report_key)
    if not PG:
        raise KeyError(
            f"reports.json の '{report_key}' に page_geometry がありません。")

    LEFT, WIDTH = K.LEFT, K.WIDTH
    BODY_TOP = PG["body_top"]
    SUMM_TOP = K.SUMM_TOP
    TWO = PG["two_col"]
    CARDS = PG["kpi_cards"]
    CWB = PG["chart_with_side_box"]
    THREE = PG["three_col"]
    CAPR = PG["capture_right"]
    CAPL = PG["capture_left"]

    def _finish(s, summ, note):
        if summ:
            K.summary(s, summ)
        if note:
            K.note(s, note)
        K.equalize()
        return s

    def _caption(s, x, y, w, h, text):
        """キャプチャの取得元URL。画像とは重ねず、外側に置く。"""
        if not text:
            return
        K.run(K.tb(s, x, y, w, h).text_frame.paragraphs[0],
              text, K.LGRAY, False, 7.5)

    def _img(s, path, x, y, max_w, max_h):
        """アスペクト比を保ったまま枠に収め、中央に置く。"""
        from PIL import Image
        Image.MAX_IMAGE_PIXELS = None
        with Image.open(path) as im:
            pw, ph = im.size
        r = min(max_w / (pw / 37.8), max_h / (ph / 37.8))
        w, h = (pw / 37.8) * r, (ph / 37.8) * r
        return s.shapes.add_picture(path, Cm(x + (max_w - w) / 2), Cm(y),
                                    Cm(w), Cm(h))

    # ---------------------------------------------------------------- 型A
    def two_col_tables(title, sub, left, right=None, insight=None,
                       summ=None, note=None, table_top=None,
                       colw_left=None, colw_right=None, rowh=0.52,
                       insight_head=None, fs=8.5):
        """表を左右に並べ、その下に解釈の箱を全幅で置く。

        right を省くと左の表だけを置く。
        insight は (見出し, lines) か lines。箱の高さは内容から自動決定する。
        """
        s = K.content(title, sub)
        top = BODY_TOP if table_top is None else table_top
        # 表が1つだけのときは全幅に広げる（右半分が空くのを避ける）
        lw = TWO["left_w"] if right else WIDTH
        K.table(s, TWO["left_x"], top, lw, left,
                colw=colw_left, rowh=rowh, fs=fs)
        bottom = top + rowh * len(left)
        if right:
            K.table(s, TWO["right_x"], top, TWO["right_w"], right,
                    colw=colw_right, rowh=rowh, fs=fs)
            bottom = max(bottom, top + rowh * len(right))
        if insight:
            head, lines = insight if isinstance(insight, tuple) else (insight_head, insight)
            K.box(s, LEFT, bottom + 0.45, WIDTH, None, head, lines, pad="mid")
        return _finish(s, summ, note)

    # ---------------------------------------------------------------- 型B
    def kpi_and_chart(title, sub, cards, chart_kind, cats, series,
                      side=None, summ=None, note=None, legend=True, fs=8):
        """上にKPIカード、下にグラフ＋右に解釈の箱。

        cards は [(見出し, [(本文行, 色, 太字), …]), …] を最大4件。
        """
        s = K.content(title, sub)
        # カードの枚数は可変。既定座標は4枚ぶんなので、枚数に応じて幅を割り直す。
        n = len(cards)
        if n == len(CARDS["xs"]):
            xs, cw = CARDS["xs"], CARDS["w"]
        else:
            gap = CARDS["xs"][1] - CARDS["xs"][0] - CARDS["w"]
            cw = (WIDTH - gap * (n - 1)) / n
            xs = [LEFT + i * (cw + gap) for i in range(n)]
        for i, (head, lines) in enumerate(cards):
            K.box(s, xs[i], CARDS["y"], cw, CARDS["h"],
                  head, lines, anchor=MSO_ANCHOR.TOP)
        K.chart(s, chart_kind, CWB["chart_x"], CWB["chart_y"],
                CWB["chart_w"], CWB["chart_h"], cats, series,
                legend=legend, fs=fs)
        if side:
            head, lines = side if isinstance(side, tuple) else (None, side)
            K.box(s, CWB["box_x"], CWB["chart_y"], CWB["box_w"], CWB["chart_h"],
                  head, lines, anchor=MSO_ANCHOR.TOP)
        return _finish(s, summ, note)

    def two_charts(title, sub, left_spec, right_spec, insight=None,
                   summ=None, note=None, chart_h=6.90, fs=8):
        """グラフを左右に並べ、その下に解釈の箱を全幅で置く。

        *_spec は (kind, cats, series, 凡例の有無)。
        """
        s = K.content(title, sub)
        for spec, x, w in ((left_spec, TWO["left_x"], TWO["left_w"]),
                           (right_spec, TWO["right_x"], TWO["right_w"])):
            if not spec:
                continue
            kind, cats, series = spec[0], spec[1], spec[2]
            legend = spec[3] if len(spec) > 3 else True
            K.chart(s, kind, x, BODY_TOP, w, chart_h, cats, series,
                    legend=legend, fs=fs)
        if insight:
            head, lines = insight if isinstance(insight, tuple) else (None, insight)
            K.box(s, LEFT, BODY_TOP + chart_h + 0.30, WIDTH, None,
                  head, lines, pad="mid")
        return _finish(s, summ, note)

    # ---------------------------------------------------------------- 型C
    def boxes(title, sub, blocks, summ=None, note=None):
        """箱だけのページ。blocks は上から順に置く指示のリスト。

          ("full", 見出し, lines, 高さ or None)
          ("cols", [(見出し, lines), …], 高さ or None)   … 最大3カラム
        """
        s = K.content(title, sub)
        y = BODY_TOP - 0.20
        for blk in blocks:
            if blk[0] == "full":
                _, head, lines, h = blk
                b = K.box(s, LEFT, y, WIDTH, h, head, lines, pad="mid")
                y += (h if h else Emu(b.height).cm) + 0.30
            elif blk[0] == "cols":
                _, cols, h = blk
                n = len(cols)
                if n == len(THREE["xs"]):
                    xs, w = THREE["xs"], THREE["w"]
                else:
                    gap = THREE["xs"][1] - THREE["xs"][0] - THREE["w"]
                    w = (WIDTH - gap * (n - 1)) / n
                    xs = [LEFT + i * (w + gap) for i in range(n)]
                tallest = 0
                for x, (head, lines) in zip(xs, cols):
                    b = K.box(s, x, y, w, h, head, lines, anchor=MSO_ANCHOR.TOP)
                    tallest = max(tallest, h if h else Emu(b.height).cm)
                K.equalize()
                y += tallest + 0.30
            if y > SUMM_TOP - 0.4:
                print(f"    ※ 箱が要約枠に近づいています（y={y:.2f}cm）: {title[:30]}")
        return _finish(s, summ, note)

    # ---------------------------------------------------------------- 型F
    def flow_diagram(title, sub, steps, side=None, footer=None,
                     summ=None, note=None, step_w=9.50, step_h=1.50,
                     step_gap=0.70, x=None, footer_h=None):
        """段階を縦に並べ、矢印でつなぐ図解。

        steps は [(見出し, [(本文, 色, 太字), …], 脇に出す注記 or None), …]。
        導線の各段でどれだけ落ちるかを見せるページで使う（ゴールデンルート・離脱率）。
        """
        from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
        from pptx.util import Cm as _Cm, Pt as _Pt

        s = K.content(title, sub)
        bx = (WIDTH - step_w) / 2 + LEFT if x is None else x
        y = BODY_TOP + 0.30
        centers = []
        for i, st in enumerate(steps):
            head, lines = st[0], st[1]
            aside = st[2] if len(st) > 2 else None
            K.box(s, bx, y, step_w, step_h, head, lines, anchor=MSO_ANCHOR.MIDDLE)
            centers.append(y + step_h)
            if aside:
                aw = LEFT + WIDTH - (bx + step_w) - 0.40
                K.run(K.tb(s, bx + step_w + 0.40, y + step_h / 2 - 0.40,
                           aw, step_h).text_frame.paragraphs[0],
                      aside, K.GRAY, True, 10)
            if i < len(steps) - 1:
                cx = bx + step_w / 2
                cn = s.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, _Cm(cx), _Cm(y + step_h),
                    _Cm(cx), _Cm(y + step_h + step_gap))
                cn.line.color.rgb = K.LGRAY
                cn.line.width = _Pt(1.5)
                tri = s.shapes.add_shape(
                    MSO_SHAPE.ISOSCELES_TRIANGLE, _Cm(cx - 0.18),
                    _Cm(y + step_h + step_gap - 0.26), _Cm(0.36), _Cm(0.26))
                tri.rotation = 180
                tri.fill.solid()
                tri.fill.fore_color.rgb = K.LGRAY
                tri.line.fill.background()
                tri.shadow.inherit = False
            y += step_h + step_gap
        if side:
            head, lines = side if isinstance(side, tuple) else (None, side)
            K.box(s, LEFT, BODY_TOP + 0.30, bx - LEFT - 0.40, None,
                  head, lines, anchor=MSO_ANCHOR.TOP)
        if footer:
            head, lines = footer if isinstance(footer, tuple) else (None, footer)
            K.box(s, LEFT, y + 0.10, WIDTH, footer_h, head, lines, pad="tight")
        return _finish(s, summ, note)

    # ---------------------------------------------------------------- 型D
    def capture_right(title, sub, boxes_left, image, caption=None,
                      summ=None, note=None, img_top=None, img_h=None):
        """左に箱（1〜2段）、右にキャプチャ。課題詳細・改善提案で使う。

        boxes_left は [(見出し, lines, 高さ), …]。
        """
        s = K.content(title, sub)
        y = BODY_TOP
        for head, lines, h in boxes_left:
            K.box(s, CAPR["box_x"], y, CAPR["box_w"], h, head, lines,
                  anchor=MSO_ANCHOR.TOP)
            y += h + 0.30
        iy = CAPR["img_y"] if img_top is None else img_top
        ih = (SUMM_TOP - 0.35 - iy - CAPR["caption_h"]) if img_h is None else img_h
        if image and os.path.exists(image):
            pic = _img(s, image, CAPR["img_x"], iy, CAPR["img_w"], ih)
            cy = iy + Emu(pic.height).cm + 0.10
        else:
            cy = iy + ih + 0.10
        _caption(s, CAPR["img_x"], cy, CAPR["img_w"], CAPR["caption_h"], caption)
        return _finish(s, summ, note)

    def capture_left(title, sub, image, boxes_right, caption=None,
                     summ=None, note=None):
        """左にキャプチャ、右に箱2段。競合比較で使う。

        boxes_right は [(見出し, lines), (見出し, lines)]（最大2段）。
        """
        s = K.content(title, sub)
        img_h = CAPL["box_top2"] + CAPL["box_h2"] - CAPL["img_y"] - CAPL["caption_h"] - 0.2
        img_w = CAPL["img_w"]
        cy = CAPL["img_y"] + img_h + 0.10
        if image and os.path.exists(image):
            pic = _img(s, image, CAPL["img_x"], CAPL["img_y"], img_w, img_h)
            cy = CAPL["img_y"] + Emu(pic.height).cm + 0.10
            # 縦長すぎる画像は枠に収めると細くなり、何が写っているか読めない。
            # 自動では切りどころを判断できないため、切り出しを促す警告だけ出す。
            if Emu(pic.width).cm < 5.0:
                print("    ※ キャプチャが細すぎます（幅 "
                      f"{Emu(pic.width).cm:.1f}cm）: "
                      f"{os.path.basename(image)}")
                print("      解説で触れた箇所を切り出してから貼ってください。")
        _caption(s, CAPL["img_x"], cy, img_w, CAPL["caption_h"], caption)
        tops = ((CAPL["box_top1"], CAPL["box_h1"]), (CAPL["box_top2"], CAPL["box_h2"]))
        for (head, lines), (ty, th) in zip(boxes_right[:2], tops):
            K.box(s, CAPL["box_x"], ty, CAPL["box_w"], th, head, lines,
                  anchor=MSO_ANCHOR.TOP)
        return _finish(s, summ, note)

    # ---------------------------------------------------------------- 競合章
    def competitor_pages(company_pages, *, conclusion=None, listing=None,
                         crosscut=None):
        """競合章をまとめて作る。1社1枚が既定（COMPETITOR_SELECTION.md）。

        company_pages は capture_left に渡す辞書のリスト。
        conclusion / listing / crosscut は、それぞれ先頭2枚と末尾1枚の指示。
        """
        made = []
        if conclusion:
            made.append(two_col_tables(**conclusion))
        if listing:
            made.append(two_col_tables(**listing))
        for spec in company_pages:
            made.append(capture_left(**spec))
        if crosscut:
            made.append(two_col_tables(**crosscut))
        return made

    # ---------------------------------------------------------------- 用語集
    def glossary(title, terms, per_page=None, sub=None, summ=None, note=None):
        """用語集。2カラムに分け、量に応じて自動でページ分割する。

        terms は [(用語, 説明), …]。
        """
        half = per_page or 18
        pages = [terms[i:i + half] for i in range(0, len(terms), half)] or [[]]
        made = []
        for pi, chunk in enumerate(pages, 1):
            t = title if len(pages) == 1 else f"{title}（{pi}/{len(pages)}）"
            s = K.content(t, sub or "本文で説明を省いた用語をまとめています。")
            mid = -(-len(chunk) // 2)
            cols = (chunk[:mid], chunk[mid:])
            xs = (LEFT, LEFT + WIDTH / 2 + 0.3)
            w = WIDTH / 2 - 0.3
            h = SUMM_TOP - BODY_TOP - 0.4
            for x, col in zip(xs, cols):
                if not col:
                    continue
                lines = []
                for term, desc in col:
                    lines.append(("■ " + term, K.NAVY, True))
                    lines.append(("　" + desc, K.GRAY, False))
                sz = K.fit_size(lines, w, h)
                K.box(s, x, BODY_TOP - 0.10, w, h, None, lines,
                      size=sz, hsize=sz + 0.5, anchor=MSO_ANCHOR.TOP)
            made.append(_finish(s, summ if pi == 1 else None,
                                note if pi == 1 else None))
        return made

    return SimpleNamespace(
        two_col_tables=two_col_tables,
        kpi_and_chart=kpi_and_chart,
        two_charts=two_charts,
        boxes=boxes,
        flow_diagram=flow_diagram,
        capture_right=capture_right,
        capture_left=capture_left,
        competitor_pages=competitor_pages,
        glossary=glossary,
        heatmap=K.heatmap_slide,
        geometry=PG,
    )
