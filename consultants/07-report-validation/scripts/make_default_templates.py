# -*- coding: utf-8 -*-
"""既定テンプレート（16:9 / 4:3）を生成する。

クライアントからテンプレートの提供が無い場合に使う、装飾を持たない素の雛形。
`TEMPLATE_POLICY.md` の §3 で用紙比率を確定させてから、対応するほうを使う。

    python scripts/make_default_templates.py [出力先フォルダ]

生成物はリポジトリにコミットせず、必要なときに各自の環境で作る。
（バイナリを公開リポジトリに置かない方針のため）
"""
from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.util import Cm, Pt

SIZES = {
    "16x9": (33.87, 19.05),
    "4x3": (25.40, 19.05),
}
# 既定の配色。ブランド指定があれば branding.json 側で上書きする。
COLORS = {
    "navy": "1F3A5F", "blue": "2E507E", "gray": "333333",
    "lgray": "808080", "light": "EEF3FA", "white": "FFFFFF",
    "red": "C00000", "green": "1B7F4B", "orange": "C16A00",
}


def build(name: str, w_cm: float, h_cm: float, outdir: str) -> str:
    prs = Presentation()
    prs.slide_width = Cm(w_cm)
    prs.slide_height = Cm(h_cm)

    # 既定テンプレートは python-pptx の標準レイアウトをそのまま使う。
    # 名前は branding.json の layouts に書き写す。
    layouts = [lo.name for lo in prs.slide_masters[0].slide_layouts]

    # 使い方を1枚だけ入れておく（生成時に deck_kit が全削除するので残らない）
    s = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])
    if s.shapes.title:
        s.shapes.title.text = f"既定テンプレート（{name.replace('x', ':')}）"
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = ("クライアント提供のテンプレートが無い場合の雛形です。\n"
                       "配色・フォントは branding.json で指定してください。")
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(14)

    os.makedirs(outdir, exist_ok=True)
    path = os.path.join(outdir, f"default_{name}.pptx")
    prs.save(path)
    return path, layouts


def main() -> int:
    sys.stdout.reconfigure(encoding="utf-8")
    outdir = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "templates")
    print("既定テンプレートを生成します。")
    for name, (w, h) in SIZES.items():
        path, layouts = build(name, w, h, outdir)
        print(f"\n  {name.replace('x', ':')}  {w} × {h} cm")
        print(f"    {path}")
        print(f"    レイアウト名（branding.json の layouts に書き写す）:")
        for i, n in enumerate(layouts[:4]):
            print(f"      [{i}] {n}")
    print("\n■ branding.json の書き方（例）")
    print('  "template":      { "path": "…/default_16x9.pptx" },')
    print('  "slide_size_cm": { "width": 33.87, "height": 19.05 },')
    print('  "layouts":       { "title_slide": "Title Slide",')
    print('                     "content": "Title and Content",')
    print('                     "section_header": "Section Header" }')
    print("\n■ 既定の配色（必要なら上書きしてください）")
    for k, v in COLORS.items():
        print(f'    "{k}": "{v}"')
    print("\n※ 生成物はコミットしないでください（公開リポジトリにバイナリを置かない方針）。")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
